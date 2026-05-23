"""
Gera a apresentação PowerPoint do Projeto 2 - SNS24 Chatbot
Corre: python gerar_apresentacao.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── PALETA ───────────────────────────────────────────────────────────────────
AZUL_ESCURO  = RGBColor(0x1A, 0x37, 0x5C)   # fundo slides título
AZUL_MED     = RGBColor(0x1E, 0x88, 0xE5)   # destaques
AZUL_CLARO   = RGBColor(0xE3, 0xF2, 0xFD)   # fundo conteúdo alternativo
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TEXTO  = RGBColor(0x21, 0x21, 0x21)
CINZA_CLARO  = RGBColor(0xF5, 0xF5, 0xF5)
VERDE        = RGBColor(0x2E, 0x7D, 0x32)
VERMELHO     = RGBColor(0xC6, 0x28, 0x28)
LARANJA      = RGBColor(0xE6, 0x5B, 0x00)
AMARELO_E    = RGBColor(0xF5, 0x7F, 0x17)
CINZA_URG    = RGBColor(0x37, 0x47, 0x4F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completamente em branco


# ── HELPERS ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=CINZA_TEXTO,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=CINZA_TEXTO,
             align=PP_ALIGN.LEFT, italic=False, space_before=0, indent=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def slide_titulo(prs, titulo, subtitulo="", autores=None):
    slide = prs.slides.add_slide(BLANK)
    # fundo gradiente simulado com dois retângulos
    add_rect(slide, 0, 0, 13.33, 7.5, AZUL_ESCURO)
    add_rect(slide, 0, 0, 13.33, 0.06, AZUL_MED)   # linha topo
    add_rect(slide, 0, 7.44, 13.33, 0.06, AZUL_MED) # linha base

    # barra decorativa lateral
    add_rect(slide, 0, 0, 0.4, 7.5, AZUL_MED)

    add_text(slide, "Universidade do Minho · Sistemas Inteligentes",
             0.7, 0.35, 11, 0.4, 13, False, RGBColor(0xBB, 0xDE, 0xFB))

    add_text(slide, titulo, 0.7, 1.2, 11.5, 2.2, 44, True, BRANCO,
             PP_ALIGN.LEFT)

    if subtitulo:
        add_text(slide, subtitulo, 0.7, 3.3, 11, 0.7, 22, False,
                 RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.LEFT)

    if autores:
        add_text(slide, autores, 0.7, 4.4, 11, 1.2, 16, False,
                 RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.LEFT)

    add_text(slide, "Maio 2026", 0.7, 6.8, 4, 0.4, 13, False,
             RGBColor(0x90, 0xCA, 0xF9))
    return slide


def slide_secao(prs, titulo_secao, num=""):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, AZUL_MED)
    add_rect(slide, 0, 0, 0.5, 7.5, AZUL_ESCURO)
    if num:
        add_text(slide, num, 1.0, 1.8, 2, 1.5, 80, True,
                 RGBColor(0xFF, 0xFF, 0xFF, ), PP_ALIGN.LEFT)
        add_text(slide, titulo_secao, 1.0, 3.2, 11, 1.2, 36, True,
                 BRANCO, PP_ALIGN.LEFT)
    else:
        add_text(slide, titulo_secao, 1.0, 2.8, 11, 1.5, 40, True,
                 BRANCO, PP_ALIGN.LEFT)
    return slide


def slide_conteudo(prs, titulo, bullets, notas=""):
    slide = prs.slides.add_slide(BLANK)
    # cabeçalho
    add_rect(slide, 0, 0, 13.33, 1.1, AZUL_ESCURO)
    add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MED)
    add_text(slide, titulo, 0.35, 0.12, 12.5, 0.85, 28, True, BRANCO)
    # fundo conteúdo
    add_rect(slide, 0, 1.15, 13.33, 6.35, CINZA_CLARO)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35),
                                     Inches(12.3), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if b.get("type") == "title":
            if not first:
                add_para(tf, "", 6)
            add_para(tf, b["text"], 20, True, AZUL_ESCURO,
                     space_before=8 if not first else 0)
        elif b.get("type") == "sub":
            add_para(tf, "    " + b["text"], 15, False, CINZA_TEXTO,
                     space_before=2, indent=1)
        elif b.get("type") == "code":
            add_para(tf, b["text"], 13, False, RGBColor(0x1A, 0x23, 0x7E),
                     italic=True, space_before=2)
        elif b.get("type") == "highlight":
            add_para(tf, b["text"], 16, True, AZUL_MED, space_before=4)
        else:
            add_para(tf, "• " + b["text"], 16, False, CINZA_TEXTO,
                     space_before=3)
        first = False

    if notas:
        add_text(slide, notas, 0.5, 7.1, 12, 0.35, 11, False,
                 RGBColor(0x75, 0x75, 0x75), italic=True)
    return slide


def slide_dois_col(prs, titulo, col1_titulo, col1_items,
                   col2_titulo, col2_items):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, AZUL_ESCURO)
    add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MED)
    add_text(slide, titulo, 0.35, 0.12, 12.5, 0.85, 28, True, BRANCO)
    add_rect(slide, 0, 1.15, 13.33, 6.35, CINZA_CLARO)

    # coluna esquerda
    add_rect(slide, 0.4, 1.4, 5.8, 5.8, BRANCO)
    add_text(slide, col1_titulo, 0.55, 1.5, 5.5, 0.5, 17, True, AZUL_ESCURO)
    txL = slide.shapes.add_textbox(Inches(0.55), Inches(2.1),
                                    Inches(5.5), Inches(4.9))
    tf = txL.text_frame
    tf.word_wrap = True
    for item in col1_items:
        add_para(tf, "• " + item, 15, False, CINZA_TEXTO, space_before=4)

    # coluna direita
    add_rect(slide, 6.9, 1.4, 6.0, 5.8, BRANCO)
    add_text(slide, col2_titulo, 7.05, 1.5, 5.7, 0.5, 17, True, AZUL_ESCURO)
    txR = slide.shapes.add_textbox(Inches(7.05), Inches(2.1),
                                    Inches(5.7), Inches(4.9))
    tf = txR.text_frame
    tf.word_wrap = True
    for item in col2_items:
        add_para(tf, "• " + item, 15, False, CINZA_TEXTO, space_before=4)
    return slide


def slide_tabela_triagem(prs):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, AZUL_ESCURO)
    add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MED)
    add_text(slide, "Base de Conhecimento — Novos Níveis de Triagem",
             0.35, 0.12, 12.5, 0.85, 28, True, BRANCO)
    add_rect(slide, 0, 1.15, 13.33, 6.35, CINZA_CLARO)

    niveis = [
        ("EMERGÊNCIA",     VERMELHO,          "Ligue 112 imediatamente",
         "r_em8–r_em11: pneumonia grave, meningite, meningococcemia, anafilaxia"),
        ("MUITO URGENTE",  RGBColor(0xB7,0x1C,0x1C), "Urgência hospitalar agora",
         "r_mu10–r_mu13: meningite/encefalite, anafilaxia, HSA, AVC-visão"),
        ("URGENTE",        LARANJA,           "Centro de saúde em 24h",
         "r_ur10–r_ur14: vómitos isolados, dor garganta+febre, bebé+vómitos, diarreia, visão"),
        ("POUCO URGENTE",  RGBColor(0xF9,0xA8,0x25), "Autocuidados, vigiar",
         "r_pu5–r_pu7: anti-sobreavaliação de quadros virais benignos"),
        ("SEM ALARME",     VERDE,             "Vigiar em casa",
         "Lógica de negação: sem sintomas de alarme confirmados"),
    ]
    y = 1.35
    for nome, cor, acao, regras in niveis:
        add_rect(slide, 0.4, y, 2.5, 0.88, cor)
        add_text(slide, nome, 0.45, y+0.18, 2.4, 0.55, 13, True, BRANCO,
                 PP_ALIGN.CENTER)
        add_rect(slide, 2.95, y, 3.6, 0.88, BRANCO)
        add_text(slide, acao, 3.0, y+0.18, 3.5, 0.55, 13, False, CINZA_TEXTO)
        add_rect(slide, 6.6, y, 6.3, 0.88, RGBColor(0xEE,0xEE,0xEE))
        add_text(slide, regras, 6.65, y+0.12, 6.2, 0.65, 12, False,
                 RGBColor(0x42,0x42,0x42), italic=True)
        y += 0.98
    return slide


def slide_arquitetura(prs):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, AZUL_ESCURO)
    add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MED)
    add_text(slide, "Arquitetura do Sistema", 0.35, 0.12, 12.5, 0.85,
             28, True, BRANCO)
    add_rect(slide, 0, 1.15, 13.33, 6.35, CINZA_CLARO)

    componentes = [
        ("Frontend\nchatbot.html", 0.4,  2.5, AZUL_MED),
        ("Backend Python\nFastAPI :8081", 3.5, 2.5, AZUL_ESCURO),
        ("LLM Ollama\nllama3.2:3b :11434", 6.6, 1.5, RGBColor(0x4A,0x14,0x8C)),
        ("Motor MYCIN\nProlog :8080",      6.6, 4.0, VERDE),
        ("Base Conhecimento\n.pl files",   10.0, 2.5, RGBColor(0x37,0x47,0x4F)),
    ]
    for label, x, y, cor in componentes:
        add_rect(slide, x, y, 2.7, 1.1, cor)
        add_text(slide, label, x+0.1, y+0.1, 2.5, 0.9, 14, True, BRANCO,
                 PP_ALIGN.CENTER)

    # setas (representadas por barras finas)
    add_rect(slide, 3.1, 2.95, 0.4, 0.12, AZUL_MED)           # front→back
    add_rect(slide, 6.2, 2.1,  0.4, 0.12, RGBColor(0x4A,0x14,0x8C)) # back→llm
    add_rect(slide, 6.2, 4.55, 0.4, 0.12, VERDE)               # back→prolog
    add_rect(slide, 9.3, 2.95, 0.7, 0.12, RGBColor(0x37,0x47,0x4F)) # prolog→kb

    # labels setas
    add_text(slide, "HTTP", 3.0, 2.75, 0.7, 0.25, 10, False,
             RGBColor(0x75,0x75,0x75), PP_ALIGN.CENTER)
    add_text(slide, "Linguagem\nNatural", 5.9, 1.85, 0.8, 0.5, 9, False,
             RGBColor(0x75,0x75,0x75), PP_ALIGN.CENTER)
    add_text(slide, "Inferência\nMYCIN", 5.9, 4.3, 0.8, 0.5, 9, False,
             RGBColor(0x75,0x75,0x75), PP_ALIGN.CENTER)
    add_text(slide, "Lê regras\nProlog", 9.1, 2.7, 1.0, 0.5, 9, False,
             RGBColor(0x75,0x75,0x75), PP_ALIGN.CENTER)

    # nota
    add_text(slide,
             "LLM → apenas geração de linguagem empática  |  "
             "Motor MYCIN → inferência clínica determinística  |  "
             "Python MYCIN → alternativa se Prolog indisponível",
             0.4, 6.9, 12.5, 0.4, 11, False,
             RGBColor(0x75,0x75,0x75), PP_ALIGN.CENTER, italic=True)
    return slide


# ── SLIDES ───────────────────────────────────────────────────────────────────

# 1. TÍTULO
slide_titulo(
    prs,
    "Projeto 2 — SNS24 Chatbot\ncom LLM e RAG",
    "Sistema Inteligente de Triagem Clínica em Linguagem Natural",
    "Alexandr Tchikoulaev  ·  André Pinto  ·  João Alves  ·  Rui Silva\n"
    "Grupo 4 · Sistemas Inteligentes de Apoio à Decisão · UMinho"
)

# 2. AGENDA
slide_conteudo(prs, "Agenda", [
    {"text": "Enquadramento — do P1 para o P2",        "type": "title"},
    {"text": "Evolução da Base de Conhecimento"},
    {"text": "Arquitetura do Sistema de Chatbot"},
    {"text": "Módulo RAG (Retrieval-Augmented Generation)", "type": "title"},
    {"text": "Deteção de Sintomas em Linguagem Natural"},
    {"text": "Motor de Inferência MYCIN em Python"},
    {"text": "Gestão da Conversa e Seleção de Perguntas", "type": "title"},
    {"text": "Integração com LLM (Llama 3.2:3b)"},
    {"text": "Interface Web"},
    {"text": "Funcionalidades Avançadas e Conclusões",   "type": "title"},
])

# 3. SECÇÃO — ENQUADRAMENTO
slide_secao(prs, "Enquadramento", "01")

slide_conteudo(prs, "Do Projeto 1 para o Projeto 2", [
    {"text": "Projeto 1 — O que foi construído", "type": "title"},
    {"text": "Motor MYCIN em Prolog com 36 regras clínicas e 28 sintomas"},
    {"text": "Pipeline de aprendizagem automática (gerar_dataset → treinar → base_conhecimento_b.pl)"},
    {"text": "Interface web (server.pl + interface.html) com triagem estruturada"},
    {"text": "Projeto 2 — O que acrescenta", "type": "title"},
    {"text": "Chatbot conversacional: o utilizador descreve os sintomas em linguagem natural"},
    {"text": "LLM local (Llama 3.2:3b via Ollama) para geração de respostas empáticas"},
    {"text": "RAG: toda a base de conhecimento Prolog indexada automaticamente"},
    {"text": "Motor de diálogo dinâmico — seleciona a pergunta mais informativa em cada turno"},
    {"text": "A base de conhecimento foi também expandida com novos sintomas e regras", "type": "highlight"},
], "O P1 é reutilizado integralmente: base Prolog, motor MYCIN, pipeline ML e triagens.csv")

# 4. SECÇÃO — BASE DE CONHECIMENTO
slide_secao(prs, "Evolução da Base de Conhecimento", "02")

slide_conteudo(prs, "Novos Sintomas Adicionados (5)", [
    {"text": "Motivação: 4 síndromes de alto risco sem representação na versão P1", "type": "highlight"},
    {"text": "rigidez_nuca", "type": "title"},
    {"text": "Rigidez da nuca → sinal clássico de meningite bacteriana"},
    {"text": "rash_petequial", "type": "title"},
    {"text": "Manchas que não desaparecem à pressão → meningococcemia"},
    {"text": "reacao_alergica_grave", "type": "title"},
    {"text": "Inchaço de face/língua/garganta → anafilaxia"},
    {"text": "dor_cabeca_subita", "type": "title"},
    {"text": '"Pior dor de cabeça da vida" → hemorragia subaracnoideia'},
    {"text": "visao_alterada", "type": "title"},
    {"text": "Alteração visual súbita → AVC (sinal BE-FAST)"},
])

slide_tabela_triagem(prs)

slide_conteudo(prs, "Decisão Arquitetural — Remoção das Contra-Evidências", [
    {"text": "Versão P1: 13 regras com CF negativo (r_c_em1 a r_c_ur3)", "type": "title"},
    {"text": "Ex.: presença de dor_leve → CF emergência −0.70"},
    {"text": "Problema crítico de segurança clínica", "type": "highlight"},
    {"text": "Um doente com paragem respiratória (CF 0.95) + mal_estar (CF −0.45) → emergência atenuada"},
    {"text": "Comportamento clinicamente inaceitável"},
    {"text": "Solução: substituição por premissas negadas", "type": "title"},
    {"text": "r_pu5: constipacao + dor_garganta + nao(febre_alta) → pouco_urgente (CF 0.72)", "type": "code"},
    {"text": "r_pu6: febre_baixa + mal_estar + nao(dor_abd) → pouco_urgente (CF 0.68)",       "type": "code"},
    {"text": "Raciocínio positivo e controlado: só dispara para nível benigno quando ausência de alarme confirmada"},
], "28 → 36 regras  ·  13 contra-evidências removidas  ·  2 regras modificadas (r_ur5, r_ur6)")

# 5. SECÇÃO — ARQUITETURA
slide_secao(prs, "Arquitetura do Sistema", "03")

slide_arquitetura(prs)

# 6. SECÇÃO — RAG
slide_secao(prs, "Módulo RAG", "04")

slide_conteudo(prs, "RAG — Retrieval-Augmented Generation", [
    {"text": "O que é o RAG neste sistema", "type": "title"},
    {"text": "No arranque: parser lê base_conhecimento_a.pl e _b.pl via regex"},
    {"text": "Constrói automaticamente 60–80 documentos em linguagem natural"},
    {"text": "3 grupos de documentos", "type": "title"},
    {"text": "Protocolo geral SNS24 (5 níveis e respetivas ações)"},
    {"text": "Uma entrada por regra de produção (premissas traduzidas para PT)"},
    {"text": "Grupos de sintomas por nível de gravidade com explicações clínicas"},
    {"text": "Motor de recuperação", "type": "title"},
    {"text": "TF-IDF com n-gramas de caráter (char_wb, bigramas–tetragramas)"},
    {"text": "Robusto a erros ortográficos e abreviaturas em português clínico"},
    {"text": "Similaridade de cosseno → top-5 documentos por mensagem"},
    {"text": "Contexto injetado nos prompts do LLM"},
])

# 7. SECÇÃO — DETEÇÃO DE SINTOMAS
slide_secao(prs, "Deteção de Sintomas", "05")

slide_dois_col(prs,
    "Deteção de Sintomas — Keywords vs LLM",
    "Sistema de Keywords (usado)",
    [
        "28 sintomas com múltiplas expressões PT",
        '"dor no peito", "aperto no peito", "dor torácica"...',
        "Deteção de negações por janela de 35 carateres",
        '"não", "sem", "nunca", "ausência de"',
        "Exclusões mútuas: febre_alta ↔ febre_baixa",
        "Determinístico — latência nula",
        "Robusto a variações ortográficas",
    ],
    "LLM para extração (desativado)",
    [
        "Função extrair_sintomas_llm existe mas não é chamada",
        "Llama 3.2:3b gerava sintomas não presentes no texto",
        "Confundia negações",
        "Devolvia JSON malformado",
        "Alucinações inaceitáveis em domínio clínico",
        "Decisão: LLM apenas para linguagem, nunca para inferência",
    ]
)

# 8. SECÇÃO — MYCIN PYTHON
slide_secao(prs, "Motor MYCIN em Python", "06")

slide_conteudo(prs, "Motor MYCIN Python — Alternativa ao Servidor Prolog", [
    {"text": "Porquê implementar o MYCIN em Python?", "type": "title"},
    {"text": "O chatbot funciona mesmo quando o servidor Prolog (P1) não está disponível"},
    {"text": "Fórmula de combinação de CFs (idêntica ao MYCIN original)", "type": "title"},
    {"text": "CF_combinado = CF₁ + CF₂ × (1 − CF₁)  [evidências positivas]",       "type": "code"},
    {"text": "Garante que a certeza cresce progressivamente sem ultrapassar 1.0"},
    {"text": "Salvaguarda de segurança clínica (sempre ativa)", "type": "title"},
    {"text": "Python e Prolog correm em paralelo quando ambos disponíveis"},
    {"text": "Se Python concluir nível MAIS GRAVE → resultado Python prevalece"},
    {"text": "Corrige: Prolog selecionava a regra com CF mais alto, não o nível mais grave",
     "type": "sub"},
    {"text": "Ex.: urgente (CF 0.85) não deve suprimir muito_urgente (CF 0.60)",
     "type": "sub"},
])

# 9. SECÇÃO — GESTÃO CONVERSA
slide_secao(prs, "Gestão da Conversa", "07")

slide_conteudo(prs, "Algoritmo de Seleção da Próxima Pergunta", [
    {"text": "Não segue uma ordem fixa — escolha dinâmica e informativa", "type": "highlight"},
    {"text": "Score por sintoma desconhecido (para cada regra ativa)", "type": "title"},
    {"text": "score = CF_regra × 2^(n_confirmadas)",                          "type": "code"},
    {"text": "Peso exponencial: favorece completar regras com evidências parciais"},
    {"text": "Regras sem nenhuma confirmação têm peso reduzido (× 0.2)",
     "type": "sub"},
    {"text": "Evita interromper uma regra em curso por uma regra genérica com CF alto",
     "type": "sub"},
    {"text": "Exemplo", "type": "title"},
    {"text": "Utilizador confirmou dor_peito → dor_irradia sobe drasticamente em score"},
    {"text": "r_mu1: dor_peito + dor_irradia → muito_urgente (CF 0.90)",        "type": "code"},
    {"text": "score(dor_irradia) = 0.90 × 2¹ = 1.80  vs  score(febre_alta) = 0.50 × 2⁰ × 0.2 = 0.10"},
])

slide_conteudo(prs, "Critérios de Disparo da Triagem e Respostas Ambíguas", [
    {"text": "Quando a triagem é desencadeada", "type": "title"},
    {"text": "Sintoma de emergência individual confirmado → imediato"},
    {"text": "Regra de emergência completamente satisfeita → early-exit absoluto (sem mais perguntas)"},
    {"text": "≥ 6 trocas e ≥ 4 sintomas avaliados (limite superior)"},
    {"text": "≥ 3 trocas, ≥ 5 sintomas e nenhuma regra de emergência ainda aberta"},
    {"text": "Tratamento de respostas ambíguas", "type": "title"},
    {"text": "Intensificadores: 'definitivamente', 'absolutamente', 'sem dúvida' → True (sem LLM)"},
    {"text": 'Qualificadas: "tenho febre, mas não acima de 39" → None (incerto)'},
    {"text": "Sintoma incerto: não é registado mas é marcado como perguntado"},
    {"text": "Intercepção de pedidos de explicação", "type": "title"},
    {"text": '"porquê precisas de saber?" → mostra explicacao(sintoma) e repete a pergunta'},
    {"text": "Estado de sessão inalterado: última pergunta mantém-se"},
])

# 10. SECÇÃO — LLM
slide_secao(prs, "Integração com LLM", "08")

slide_conteudo(prs, "Llama 3.2:3b — Papel Exclusivamente Linguístico", [
    {"text": "Configuração", "type": "title"},
    {"text": "Modelo: llama3.2:3b via Ollama (execução local, porta 11434)"},
    {"text": "Temperatura: 0.1 → mínima variabilidade, redução de alucinações"},
    {"text": "Tokens máximos: 25 (acknowledgments) / 35 (transições)"},
    {"text": "3 tipos de prompts", "type": "title"},
    {"text": "Mid-conversation: frase empática de transição antes de cada pergunta clínica"},
    {"text": "Introdução ao resultado: 1 frase antes do diagnóstico (sem revelar nível)"},
    {"text": "Português europeu formal: 'si' em vez de 'você' (corrigido por regex)"},
    {"text": "Pós-processamento obrigatório", "type": "title"},
    {"text": "Remove '112', 'urgência', 'emergência', 'grave' do texto intercalar"},
    {"text": "Deteta antecipação do próximo sintoma → substitui por texto neutro fixo"},
    {"text": "Trunca a 2 frases máximo"},
    {"text": "LLM nunca influencia a lógica de triagem", "type": "highlight"},
])

# 11. SECÇÃO — INTERFACE
slide_secao(prs, "Interface Web", "09")

slide_conteudo(prs, "chatbot.html — Single Page Application", [
    {"text": "Tecnologias: HTML + CSS + JavaScript puro (sem framework externo)", "type": "title"},
    {"text": "Painel de chat (esquerda)", "type": "title"},
    {"text": "Conversa em linguagem natural com o utilizador"},
    {"text": "Mensagens do sistema com frases empáticas geradas pelo LLM"},
    {"text": "Resultado final com nível codificado por cor e justificação clínica"},
    {"text": "Painel de diagnóstico (direita)", "type": "title"},
    {"text": "Lista em tempo real dos sintomas detetados (✓ presente / ✗ ausente)"},
    {"text": "Status de ligação ao modelo Ollama"},
    {"text": "Endpoints utilizados", "type": "title"},
    {"text": "POST /api/chat/start  → UUID de sessão + mensagem de boas-vindas",  "type": "code"},
    {"text": "POST /api/chat/message → processa mensagem + devolve resposta + estado", "type": "code"},
    {"text": "GET  /api/chat/status  → confirma llama3.2:3b carregado",           "type": "code"},
])

# 12. SECÇÃO — FUNCIONALIDADES AVANÇADAS
slide_secao(prs, "Funcionalidades Avançadas", "10")

slide_conteudo(prs, "Funcionalidades Avançadas e Persistência", [
    {"text": "Persistência automática de triagens (salvar_triagem_csv)", "type": "title"},
    {"text": "No P1: persistência dependia do estado Prolog → inoperante em modo fallback"},
    {"text": "No P2: Python persiste sempre, independentemente do motor usado"},
    {"text": "Converte session['sintomas'] → linha binária em triagens.csv"},
    {"text": "Fecha o ciclo de aprendizagem contínua do P1"},
    {"text": "Segurança clínica dupla", "type": "title"},
    {"text": "Prolog + Python correm em paralelo → resultado mais grave prevalece"},
    {"text": "Early-exit absoluto para emergências confirmadas"},
    {"text": "Explicações clínicas em contexto", "type": "title"},
    {"text": '"porquê?" → mostra explicacao/2 da base Prolog e repete a pergunta'},
    {"text": "Sem modificar o estado da sessão"},
], "Integração total com o P1: base Prolog reutilizada · triagens.csv partilhado · pipeline ML inalterado")

# 13. SECÇÃO — CONCLUSÕES
slide_secao(prs, "Conclusões", "11")

slide_conteudo(prs, "Síntese e Autoavaliação", [
    {"text": "O que foi entregue", "type": "title"},
    {"text": "Chatbot conversacional SNS24 com LLM + RAG + MYCIN"},
    {"text": "Base de conhecimento expandida: 28 → 36 regras, 5 novos sintomas"},
    {"text": "Motor MYCIN em Python (alternativa ao Prolog)"},
    {"text": "Salvaguarda de segurança clínica (nivel mais grave prevalece)"},
    {"text": "Persistência de triagens independente do motor utilizado"},
    {"text": "Principal lição aprendida", "type": "title"},
    {"text": "LLMs de 3B parâmetros são inadequados para extração estruturada em domínios clínicos"},
    {"text": "Solução: LLM para linguagem natural, keywords para extração, MYCIN para inferência"},
    {"text": "Autoavaliação: 18 valores", "type": "highlight"},
    {"text": "Alexandr Tchikoulaev  ·  André Pinto  ·  João Alves  ·  Rui Silva",
     "type": "sub"},
])

# 14. SLIDE FINAL
slide_titulo(
    prs,
    "Obrigado",
    "Questões?",
    "Alexandr Tchikoulaev  ·  André Pinto  ·  João Alves  ·  Rui Silva\n"
    "Grupo 4 · Sistemas Inteligentes de Apoio à Decisão · UMinho · Maio 2026"
)

# ── GUARDAR ──────────────────────────────────────────────────────────────────
out = "apresentacao_p2.pptx"
prs.save(out)
print(f"OK - Apresentacao guardada: {out}  ({prs.slides.__len__()} slides)")
